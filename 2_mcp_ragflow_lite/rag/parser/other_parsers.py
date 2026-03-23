"""
TXT / HTML / JSON / PPT 解析器
轻量级实现，全部继承 BaseParser
"""
import logging
import re
import json as json_lib
from typing import Optional

from rag.parser.base import BaseParser
from common.registry import parser_registry

logger = logging.getLogger(__name__)


# ==================== TXT ====================

@parser_registry.register(".txt")
@parser_registry.register(".csv")
class TxtParser(BaseParser):
    """纯文本 / CSV 解析器"""

    name = "TXT Parser"
    extensions = [".txt", ".csv"]

    def parse(self, filename: str, binary: Optional[bytes] = None):
        try:
            if binary:
                from rag.nlp import find_codec
                codec = find_codec(binary)
                text = binary.decode(codec)
            else:
                with open(filename, "r", encoding="utf-8") as f:
                    text = f.read()
        except Exception as e:
            logger.error(f"Failed to read TXT {filename}: {e}")
            return [], []

        sections = []
        paragraphs = re.split(r"\n\s*\n", text)
        for para in paragraphs:
            para = para.strip()
            if para:
                sections.append((para, "text"))

        return sections, []


# ==================== HTML ====================

@parser_registry.register(".html")
@parser_registry.register(".htm")
class HtmlParser(BaseParser):
    """HTML 文档解析器"""

    name = "HTML Parser"
    extensions = [".html", ".htm"]

    def parse(self, filename: str, binary: Optional[bytes] = None):
        from bs4 import BeautifulSoup

        try:
            if binary:
                from rag.nlp import find_codec
                codec = find_codec(binary)
                html = binary.decode(codec)
            else:
                with open(filename, "r", encoding="utf-8") as f:
                    html = f.read()
        except Exception as e:
            logger.error(f"Failed to read HTML {filename}: {e}")
            return [], []

        soup = BeautifulSoup(html, "html.parser")

        # 移除 script 和 style
        for tag in soup.find_all(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        sections = []
        tables = []

        # 提取表格
        for table in soup.find_all("table"):
            tables.append(str(table))
            table.decompose()

        # 提取正文段落
        for elem in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "div"]):
            text = elem.get_text(strip=True)
            if text and len(text) > 1:
                tag = "title" if elem.name in ["h1", "h2", "h3"] else "text"
                sections.append((text, tag))

        # 如果段落提取结果很少，fallback 到全文本
        if not sections:
            text = soup.get_text()
            paragraphs = re.split(r"\n\s*\n", text)
            for para in paragraphs:
                para = para.strip()
                if para:
                    sections.append((para, "text"))

        return sections, tables


# ==================== JSON ====================

@parser_registry.register(".json")
class JsonParser(BaseParser):
    """JSON 文档解析器"""

    name = "JSON Parser"
    extensions = [".json"]

    def parse(self, filename: str, binary: Optional[bytes] = None):
        try:
            if binary:
                from rag.nlp import find_codec
                codec = find_codec(binary)
                text = binary.decode(codec)
            else:
                with open(filename, "r", encoding="utf-8") as f:
                    text = f.read()
            data = json_lib.loads(text)
        except Exception as e:
            logger.error(f"Failed to read JSON {filename}: {e}")
            return [], []

        sections = []

        def flatten(obj, prefix=""):
            if isinstance(obj, dict):
                for k, v in obj.items():
                    key = f"{prefix}.{k}" if prefix else k
                    flatten(v, key)
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    flatten(item, f"{prefix}[{i}]")
            else:
                text = f"{prefix}: {obj}" if prefix else str(obj)
                if text.strip():
                    sections.append((text, "text"))

        flatten(data)
        return sections, []


# ==================== PPT ====================

@parser_registry.register(".pptx")
@parser_registry.register(".ppt")
class PptParser(BaseParser):
    """PPT/PPTX 文档解析器"""

    name = "PPT Parser"
    extensions = [".pptx", ".ppt"]

    def parse(self, filename: str, binary: Optional[bytes] = None):
        from pptx import Presentation
        import io

        try:
            if binary:
                prs = Presentation(io.BytesIO(binary))
            else:
                prs = Presentation(filename)
        except Exception as e:
            logger.error(f"Failed to open PPT {filename}: {e}")
            return [], []

        sections = []
        tables = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        header = rows[0]
                        html = "<table>\n<tr>" + "".join(f"<th>{h}</th>" for h in header) + "</tr>\n"
                        for row in rows[1:]:
                            html += "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>\n"
                        html += "</table>"
                        tables.append(html)

            if slide_texts:
                sections.append(("\n".join(slide_texts), f"slide_{slide_num}"))

        return sections, tables


# ── 向后兼容模块级函数 ──
_txt = TxtParser()
_html = HtmlParser()
_json = JsonParser()
_ppt = PptParser()

parse_txt = _txt.parse
parse_html = _html.parse
parse_json = _json.parse
parse_ppt = _ppt.parse
